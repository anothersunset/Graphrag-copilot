"""GraphRAG Copilot - 多模态文档解析服务"""
from typing import List, Dict, Any
from pathlib import Path
import hashlib
from config.settings import settings

class DocumentParser:
    def __init__(self):
        self.supported_formats = {
            ".pdf": self.parse_pdf,
            ".docx": self.parse_docx,
            ".pptx": self.parse_pptx,
            ".txt": self.parse_txt,
            ".md": self.parse_txt,
            ".jpg": self.parse_image,
            ".jpeg": self.parse_image,
            ".png": self.parse_image,
            ".mp4": self.parse_video,
            ".wav": self.parse_audio,
            ".mp3": self.parse_audio,
        }

    def parse(self, file_path: str) -> Dict[str, Any]:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError("文件不存在: " + file_path)

        suffix = path.suffix.lower()
        if suffix not in self.supported_formats:
            raise ValueError("不支持的文件格式: " + suffix)

        file_hash = self._calculate_hash(path)
        parser = self.supported_formats[suffix]
        content = parser(path)

        return {
            "file_path": str(path),
            "file_name": path.name,
            "file_type": suffix,
            "file_hash": file_hash,
            "content": content,
            "metadata": {
                "size": path.stat().st_size,
                "modified": path.stat().st_mtime,
            },
        }

    def parse_pdf(self, path: Path) -> Dict[str, Any]:
        try:
            from pypdf2 import PdfReader
            reader = PdfReader(str(path))
            pages = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text.strip():
                    pages.append({"page": i + 1, "content": text.strip()})
            return {
                "type": "pdf",
                "total_pages": len(reader.pages),
                "pages": pages,
                "full_text": "\n\n".join([p["content"] for p in pages]),
            }
        except Exception as e:
            return {"type": "pdf", "error": str(e), "full_text": ""}

    def parse_docx(self, path: Path) -> Dict[str, Any]:
        try:
            from docx import Document
            doc = Document(str(path))
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            tables = []
            for table in doc.tables:
                tables.append([[cell.text for cell in row.cells] for row in table.rows])
            return {
                "type": "docx",
                "paragraphs": paragraphs,
                "tables": tables,
                "full_text": "\n\n".join(paragraphs),
            }
        except Exception as e:
            return {"type": "docx", "error": str(e), "full_text": ""}

    def parse_pptx(self, path: Path) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            slides = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    slides.append({"slide": i + 1, "content": "\n".join(slide_text)})
            return {
                "type": "pptx",
                "total_slides": len(prs.slides),
                "slides": slides,
                "full_text": "\n\n".join([s["content"] for s in slides]),
            }
        except Exception as e:
            return {"type": "pptx", "error": str(e), "full_text": ""}

    def parse_txt(self, path: Path) -> Dict[str, Any]:
        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()
            return {"type": "text", "content": content, "full_text": content}
        except Exception as e:
            return {"type": "text", "error": str(e), "full_text": ""}

    def parse_image(self, path: Path) -> Dict[str, Any]:
        try:
            from paddleocr import PaddleOCR
            ocr = PaddleOCR(use_angle_cls=True, lang="ch")
            result = ocr.ocr(str(path), cls=True)
            texts = []
            for line in result[0]:
                texts.append({"text": line[1][0], "confidence": line[1][1], "bbox": line[0]})
            return {
                "type": "image",
                "ocr_results": texts,
                "full_text": "\n".join([t["text"] for t in texts]),
            }
        except Exception as e:
            return {"type": "image", "error": str(e), "full_text": ""}

    def parse_video(self, path: Path) -> Dict[str, Any]:
        try:
            import whisper
            model = whisper.load_model(settings.ASR_MODEL)
            result = model.transcribe(str(path), language="zh")
            segments = [
                {"start": seg["start"], "end": seg["end"], "text": seg["text"]}
                for seg in result["segments"]
            ]
            return {
                "type": "video",
                "duration": result.get("duration", 0),
                "segments": segments,
                "full_text": result["text"],
            }
        except Exception as e:
            return {"type": "video", "error": str(e), "full_text": ""}

    def parse_audio(self, path: Path) -> Dict[str, Any]:
        return self.parse_video(path)

    def _calculate_hash(self, path: Path) -> str:
        sha256 = hashlib.sha256()
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                sha256.update(chunk)
        return sha256.hexdigest()

    def chunk_text(self, text: str, chunk_size: int = None, overlap: int = None) -> List[str]:
        chunk_size = chunk_size or settings.CHUNK_SIZE
        overlap = overlap or settings.CHUNK_OVERLAP
        if len(text) <= chunk_size:
            return [text]
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunks.append(text[start:end])
            start = end - overlap
        return chunks

doc_parser = DocumentParser()
